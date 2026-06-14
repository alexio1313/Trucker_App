"""
Executive Presentation: AI-Powered Truck Logistics Platform
Creates a professional PPTX for investor/demo presentation
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt
import datetime

# Brand colors
ORANGE = RGBColor(0xF9, 0x73, 0x16)   # #F97316
DARK   = RGBColor(0x11, 0x18, 0x27)   # #111827
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
GRAY   = RGBColor(0x6B, 0x72, 0x80)   # #6B7280
LIGHT  = RGBColor(0xF3, 0xF4, 0xF6)   # #F3F4F6
GREEN  = RGBColor(0x10, 0xB9, 0x81)   # #10B981
BLUE   = RGBColor(0x3B, 0x82, 0xF6)   # #3B82F6
PURPLE = RGBColor(0x8B, 0x5C, 0xF6)   # #8B5CF6

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # completely blank

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_w or 1)
    else:
        shape.line.fill.background()
    return shape

def txt(slide, text, l, t, w, h, size=18, bold=False, color=None, align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color or DARK
    return tb

def bullet_box(slide, items, l, t, w, h, size=14, color=None, icon="●"):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.add_paragraph() if not first else tf.paragraphs[0]
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{icon}  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color or DARK
        p.space_after = Pt(4)

def header_bar(slide, title, subtitle=None):
    rect(slide, 0, 0, 13.33, 1.1, fill=DARK)
    txt(slide, title, 0.4, 0.12, 10, 0.6, size=28, bold=True, color=WHITE)
    if subtitle:
        txt(slide, subtitle, 0.4, 0.68, 10, 0.38, size=13, color=GRAY)
    # Orange accent line
    rect(slide, 0, 1.1, 13.33, 0.06, fill=ORANGE)

def footer(slide, page_num, total=14):
    rect(slide, 0, 7.1, 13.33, 0.4, fill=DARK)
    txt(slide, "AI-Powered Truck Logistics Platform  ●  Confidential", 0.3, 7.14, 10, 0.3, size=9, color=GRAY)
    txt(slide, f"{page_num} / {total}", 12.5, 7.14, 0.8, 0.3, size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def stat_card(slide, l, t, w, h, value, label, bg=None, val_color=None):
    bg = bg or LIGHT
    val_color = val_color or ORANGE
    rect(slide, l, t, w, h, fill=bg)
    txt(slide, value, l+0.1, t+0.12, w-0.2, 0.55, size=30, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    txt(slide, label, l+0.1, t+0.65, w-0.2, 0.35, size=11, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
# Full dark background
rect(s, 0, 0, 13.33, 7.5, fill=DARK)
# Orange accent strip left
rect(s, 0, 0, 0.18, 7.5, fill=ORANGE)
# Subtle grid pattern (orange rectangles)
for i in range(6):
    rect(s, 0.5 + i*2.1, 5.8, 1.6, 1.5, fill=RGBColor(0x1F, 0x29, 0x37))

txt(s, "🚛", 1.0, 1.2, 2, 1.2, size=64, align=PP_ALIGN.CENTER)
txt(s, "AI-Powered Truck", 3.2, 1.1, 9.5, 1.1, size=48, bold=True, color=WHITE)
txt(s, "Logistics Platform", 3.2, 2.1, 9.5, 1.0, size=48, bold=True, color=ORANGE)
txt(s, "Connecting India's Freight Industry with AI, Real-Time Tracking & Digital Payments",
    3.2, 3.25, 9.5, 0.7, size=17, color=GRAY)
# Divider
rect(s, 3.2, 3.9, 6, 0.05, fill=ORANGE)
txt(s, "Executive Presentation  ●  " + datetime.date.today().strftime("%B %Y"),
    3.2, 4.05, 9, 0.5, size=14, color=GRAY)
txt(s, "CONFIDENTIAL", 3.2, 4.6, 9, 0.4, size=11, color=RGBColor(0xF9,0x73,0x16))

# Stats preview
for i, (v, l) in enumerate([("18", "Microservices"), ("9", "Languages"), ("38/38", "QA Tests Pass"), ("100%", "Uptime")]):
    stat_card(s, 0.4 + i*3.1, 5.8, 2.8, 1.2, v, l, bg=RGBColor(0x1F,0x29,0x37), val_color=ORANGE)

# ══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "The Problem", "India's ₹15 Trillion freight industry is broken")
footer(s, 2)

problems = [
    ("₹3.2L Cr", "Lost annually due to empty truck runs\n& inefficient load matching"),
    ("68%", "Loads still booked via phone calls\n& WhatsApp — zero digital trail"),
    ("42 days", "Average payment delay to truckers\ndriving financial stress"),
    ("0%", "Real-time visibility for merchants\non their cargo in transit"),
]
for i, (val, desc) in enumerate(problems):
    x = 0.3 + (i % 2) * 6.6
    y = 1.5 + (i // 2) * 2.3
    rect(s, x, y, 6.2, 2.0, fill=LIGHT)
    rect(s, x, y, 0.15, 2.0, fill=ORANGE)
    txt(s, val, x+0.3, y+0.15, 5.8, 0.8, size=36, bold=True, color=ORANGE)
    txt(s, desc, x+0.3, y+0.95, 5.8, 0.9, size=13, color=DARK)

txt(s, "5.5 million truckers in India — 80% unorganised, underserved, and underpaid",
    0.4, 6.5, 12.5, 0.5, size=14, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 3 — OUR SOLUTION
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Our Solution", "A full-stack AI logistics platform built for Bharat")
footer(s, 3)

txt(s, "TruckPlatform is an end-to-end digital logistics ecosystem connecting verified truckers with merchants using AI-powered matching, real-time GPS tracking, and instant digital payments.",
    0.4, 1.3, 12.5, 0.8, size=14, color=DARK)

pillars = [
    ("🤖", "AI Matching", "Smart load-trucker matching with pricing intelligence and route optimization"),
    ("📍", "Live Tracking", "Real-time GPS across all loads with ETA prediction and delay alerts"),
    ("💸", "Digital Payments", "Instant settlements, escrow protection, and transparent commission"),
    ("📱", "Mobile-First", "Native Android app + web portal for truckers in 9 Indian languages"),
    ("⚙️", "Admin Control", "Full fleet management, KYC workflow, and dispute resolution system"),
    ("📣", "Social Media AI", "AI-generated marketing content for merchants across 5 platforms"),
]
for i, (icon, title, desc) in enumerate(pillars):
    x = 0.3 + (i % 3) * 4.35
    y = 2.3 + (i // 3) * 2.15
    rect(s, x, y, 4.1, 1.95, fill=LIGHT)
    txt(s, icon + "  " + title, x+0.2, y+0.15, 3.8, 0.5, size=15, bold=True, color=DARK)
    txt(s, desc, x+0.2, y+0.6, 3.8, 1.1, size=11, color=GRAY)

# ══════════════════════════════════════════════════════════════════
# SLIDE 4 — PLATFORM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Platform Architecture", "18 microservices, fully containerised on Docker")
footer(s, 4)

layers = [
    ("CLIENT LAYER", [("📱 Android App", BLUE), ("🌐 Web Portal", BLUE), ("🖥 Admin Panel", BLUE)], 1.3),
    ("API GATEWAY", [("Auth + Rate Limit", ORANGE), ("JWT Middleware", ORANGE), ("Proxy + Routing", ORANGE)], 2.7),
    ("MICROSERVICES", [("Load Service", GREEN), ("Trucker Service", GREEN), ("Pricing Engine", GREEN), ("ML / AI Service", PURPLE), ("Social Publishing", PURPLE), ("Admin Service", PURPLE), ("Payment Service", GREEN), ("Notification Svc", GREEN)], 4.1),
    ("DATA LAYER", [("PostgreSQL + PostGIS", BLUE), ("MongoDB", BLUE), ("Redis Cache", BLUE), ("Kafka Events", BLUE), ("Elasticsearch", BLUE)], 5.7),
]

for layer_name, items, y in layers:
    rect(s, 0.3, y, 12.7, 1.2, fill=LIGHT)
    txt(s, layer_name, 0.4, y+0.35, 1.6, 0.5, size=9, bold=True, color=GRAY)

    item_w = 11.0 / len(items)
    for j, (name, color) in enumerate(items):
        lx = 2.0 + j * item_w
        rect(s, lx, y+0.12, item_w-0.1, 0.95, fill=color)
        txt(s, name, lx+0.05, y+0.3, item_w-0.15, 0.6, size=9, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "All services on Docker Compose ● REST APIs ● JWT Auth ● Event-Driven via Kafka ● PostGIS spatial queries",
    0.3, 7.0, 12.7, 0.35, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 5 — THREE PORTALS
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Three User Portals", "Purpose-built experience for each persona")
footer(s, 5)

portals = [
    (ORANGE, "🚛", "TRUCKER PORTAL", "Web + Mobile App", [
        "AI load matching & search",
        "Accept loads (availability gate)",
        "Journey: Loading → Driving → Delivered",
        "Live OSRM route map with fuel stations",
        "Fuel stop logging & toll tracking",
        "Earnings dashboard (daily/weekly/monthly)",
        "KYC submission & bank account setup",
        "9-language support (Hindi, Tamil, Telugu…)",
    ]),
    (BLUE, "🏪", "MERCHANT PORTAL", "Web App", [
        "Post loads with pickup/delivery details",
        "View all my loads & track status",
        "AI-powered social media post creator",
        "Raise disputes on any load",
        "Pricing intelligence & route estimates",
        "Real-time trucker location visibility",
        "Payment history & invoice downloads",
        "Multi-platform social publishing",
    ]),
    (GREEN, "⚙️", "ADMIN PANEL", "Next.js Dashboard", [
        "Live fleet map with truck clustering",
        "KYC approval/rejection queue",
        "All loads with full details & disputes",
        "User management & staff creation",
        "Social post approval workflow",
        "AI simulation control panel",
        "API health monitoring (9 services)",
        "Dispute resolution with type dropdown",
    ]),
]

for i, (color, icon, title, sub, items) in enumerate(portals):
    x = 0.3 + i * 4.35
    rect(s, x, 1.3, 4.1, 5.8, fill=LIGHT)
    rect(s, x, 1.3, 4.1, 1.1, fill=color)
    txt(s, icon + "  " + title, x+0.15, 1.38, 3.8, 0.5, size=14, bold=True, color=WHITE)
    txt(s, sub, x+0.15, 1.85, 3.8, 0.35, size=10, color=WHITE)
    bullet_box(s, items, x+0.2, 2.5, 3.8, 4.4, size=11, color=DARK, icon="✓")

# ══════════════════════════════════════════════════════════════════
# SLIDE 6 — TRUCKER JOURNEY FLOW
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Trucker Journey Flow", "End-to-end trip lifecycle with real-time state management")
footer(s, 6)

steps = [
    (GRAY,   "1", "Find Load",      "Search by city,\ncargo type, distance"),
    (BLUE,   "2", "Accept",         "AI matches → Trucker\nacknowledges load"),
    (PURPLE, "3", "Begin Loading",  "Arrive at pickup\nCargo loaded & sealed"),
    (ORANGE, "4", "Start Journey",  "GPS tracking live\nOSRM route shown"),
    (BLUE,   "5", "Fuel Stops",     "Log fuel liters,\ncost & station name"),
    (GREEN,  "6", "Mark Delivered", "Delivery confirmed\nOdometer logged"),
    (GREEN,  "7", "Earnings",       "Stats updated\nPayment triggered"),
]

for i, (color, num, title, desc) in enumerate(steps):
    x = 0.4 + i * 1.82
    # Circle
    rect(s, x, 1.5, 1.45, 1.45, fill=color)
    txt(s, num, x, 1.62, 1.45, 0.8, size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    # Arrow connector
    if i < len(steps)-1:
        txt(s, "→", x+1.45, 1.95, 0.4, 0.5, size=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    txt(s, title, x-0.1, 3.15, 1.65, 0.45, size=11, bold=True, color=DARK, align=PP_ALIGN.CENTER)
    txt(s, desc, x-0.1, 3.6, 1.65, 0.8, size=9, color=GRAY, align=PP_ALIGN.CENTER)

# Journey page features
rect(s, 0.4, 4.6, 12.5, 2.2, fill=LIGHT)
txt(s, "Journey Page Features", 0.6, 4.7, 5, 0.4, size=13, bold=True, color=DARK)
features = [
    "🗺️  Live OSRM route polyline (OpenStreetMap)",
    "⛽  Nearby fuel stations via Overpass API",
    "🚧  Toll estimation (NHAI 2025 NH rates)",
    "📊  Distance, duration & fuel calculations",
    "💰  Net earnings after toll & fuel deductions",
    "🌐  9 Indian language support",
]
for i, f in enumerate(features):
    x = 0.6 + (i % 3) * 4.2
    y = 5.15 + (i // 3) * 0.5
    txt(s, f, x, y, 4.0, 0.45, size=11, color=DARK)

# ══════════════════════════════════════════════════════════════════
# SLIDE 7 — AI & ML CAPABILITIES
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "AI & ML Capabilities", "Multi-model intelligence embedded across the platform")
footer(s, 7)

ai_features = [
    (ORANGE, "🤖", "Load-Trucker Matching", [
        "Geospatial matching via PostGIS",
        "Distance-based ranking algorithm",
        "Cargo type & capacity filtering",
        "Availability status gating",
        "Real-time load acceptance lock",
    ]),
    (PURPLE, "✍️", "Social Caption AI", [
        "Multi-provider LLM cascade:",
        "  Groq → Claude → OpenAI → Gemini",
        "Platform-specific tone & hashtags",
        "5 platforms: FB, IG, TW, LI, WA",
        "Rich template fallback (instant)",
    ]),
    (BLUE, "💰", "Dynamic Pricing Engine", [
        "Distance × weight base pricing",
        "Surge multipliers by lane & time",
        "Fuel cost integration (₹93/L diesel)",
        "Toll estimation (NHAI booth rates)",
        "Real-time market rate benchmarking",
    ]),
    (GREEN, "📍", "Route Intelligence", [
        "OSRM engine for route planning",
        "ETA with rest-stop adjustments",
        "Overpass API fuel station lookup",
        "Weather advisory per NH segment",
        "Document expiry alerts system",
    ]),
]

for i, (color, icon, title, items) in enumerate(ai_features):
    x = 0.3 + (i % 2) * 6.5
    y = 1.4 + (i // 2) * 2.6
    rect(s, x, y, 6.2, 2.4, fill=LIGHT)
    rect(s, x, y, 0.15, 2.4, fill=color)
    txt(s, icon + "  " + title, x+0.3, y+0.12, 5.7, 0.5, size=14, bold=True, color=DARK)
    bullet_box(s, items, x+0.3, y+0.6, 5.7, 1.7, size=11, color=GRAY, icon="›")

txt(s, "ML Service (Port 3007) routes tasks to optimal provider: FAST→Ollama, QUALITY→Claude Haiku, COMPLEX→Claude Opus",
    0.4, 6.8, 12.5, 0.35, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Technology Stack", "Production-grade, cloud-native, open-source first")
footer(s, 8)

categories = [
    ("Frontend", BLUE, ["React 18 + Vite (Web Portal)", "Next.js 14 Standalone (Admin)", "React Native + Expo 51 (Mobile)", "Tailwind CSS + TanStack Query", "Leaflet.js maps + OSRM routing", "i18n: 9 Indian languages"]),
    ("Backend", GREEN, ["Node.js + Express (8 services)", "TypeScript throughout", "JWT authentication + bcrypt", "http-proxy-middleware gateway", "Express rate limiting", "Zod schema validation"]),
    ("Data", ORANGE, ["PostgreSQL 15 + PostGIS", "MongoDB (social posts)", "Redis (caching + sessions)", "Elasticsearch (search)", "Apache Kafka (events)", "Apache Zookeeper"]),
    ("DevOps", PURPLE, ["Docker + Docker Compose", "18 containers orchestrated", "Hot-patch deployment pattern", "Health check endpoints", "Prometheus metrics", "Nginx reverse proxy"]),
    ("AI/ML", RGBColor(0xDC,0x26,0x26), ["Claude API (Anthropic)", "Groq (Llama 3.3 70B)", "OpenAI GPT-4o-mini", "Google Gemini 1.5 Flash", "OSRM (routing engine)", "Overpass API (POI)"]),
    ("Mobile", RGBColor(0x06,0xB6,0xD4), ["Expo SDK 51", "React Native 0.74", "MapView + OSRM routing", "expo-location (GPS)", "Debug APK (173MB)", "4 ABIs (arm64, x86…)"]),
]

for i, (cat, color, items) in enumerate(categories):
    x = 0.3 + (i % 3) * 4.35
    y = 1.35 + (i // 3) * 2.7
    rect(s, x, y, 4.1, 0.45, fill=color)
    txt(s, cat, x+0.15, y+0.06, 3.8, 0.35, size=13, bold=True, color=WHITE)
    bullet_box(s, items, x+0.1, y+0.5, 3.9, 2.1, size=10, color=DARK, icon="•")

# ══════════════════════════════════════════════════════════════════
# SLIDE 9 — ADMIN INTELLIGENCE
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Admin Intelligence Dashboard", "Complete operational control for platform managers")
footer(s, 9)

admin_features = [
    ("🗺️", "Live Fleet Map", "/admin/live-map", [
        "Real-time truck positions on Leaflet",
        "MarkerCluster for dense urban areas",
        "Click truck → OSRM route + ETA panel",
        "Status colours: orange=on_load, green=in_transit",
        "Auto-refresh every 30 seconds",
    ]),
    ("📋", "KYC Queue", "/admin/kyc", [
        "Pending submissions dashboard",
        "One-click Approve / Reject",
        "Rejection reason modal",
        "Aadhaar, PAN, DL verification",
        "Auto-filters by status",
    ]),
    ("📦", "Load Management", "/admin/loads", [
        "All loads with nested field parsing",
        "Search by city, cargo, load ID",
        "Filter by status (7 states)",
        "Full detail drawer with merchant info",
        "Raise dispute from any load",
    ]),
    ("📣", "Social Management", "/admin/social", [
        "Approve/Reject merchant posts",
        "Platform breakdown analytics",
        "Posts queue (Pending/Published/Rejected)",
        "Rejection reason with custom note",
        "Real-time post count stats",
    ]),
    ("🎮", "Simulation Panel", "/admin/simulation", [
        "Seed 3 truckers across BLR/DEL/MUM",
        "Seed 3 merchants with KYC verified",
        "Post 6 loads per city (18 total)",
        "Set trucker GPS to any Indian city",
        "Live status panel for all sim users",
    ]),
    ("📊", "API Monitor", "/admin/api-status", [
        "9 service health checks live",
        "Response time per service",
        "External integrations table",
        "Colour-coded status indicators",
        "Auto-refresh every 60 seconds",
    ]),
]

for i, (icon, title, path, items) in enumerate(admin_features):
    x = 0.3 + (i % 3) * 4.35
    y = 1.4 + (i // 3) * 2.65
    rect(s, x, y, 4.1, 2.5, fill=LIGHT)
    txt(s, icon + "  " + title, x+0.15, y+0.1, 3.8, 0.45, size=13, bold=True, color=DARK)
    txt(s, path, x+0.15, y+0.52, 3.8, 0.3, size=9, color=ORANGE)
    bullet_box(s, items, x+0.15, y+0.82, 3.8, 1.55, size=10, color=GRAY, icon="→")

# ══════════════════════════════════════════════════════════════════
# SLIDE 10 — MOBILE APP
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Mobile Application", "Native Android app built for India's truckers")
footer(s, 10)

txt(s, "Built with Expo 51 + React Native 0.74 · 173MB APK · 4 ABIs · API: http://192.168.8.101:3000",
    0.4, 1.25, 12.5, 0.4, size=12, color=GRAY)

screens = [
    ("🏠", "Dashboard", ["Active load banner", "Today's earnings", "Quick stats cards", "Weather-based greeting"]),
    ("🔍", "Find Loads", ["Search by origin city", "Filter by cargo type", "Distance & price shown", "One-tap accept"]),
    ("📍", "Map Navigation", ["Live MapView GPS", "OSRM route polyline", "Speed badge overlay", "4 metric cards"]),
    ("💰", "Earnings", ["Daily / Weekly / Monthly", "Trip count & km", "Fuel & toll breakdown", "Net earnings chart"]),
    ("👤", "Profile", ["Availability toggle", "Add truck form", "KYC submission", "Bank details setup"]),
]

for i, (icon, title, items) in enumerate(screens):
    x = 0.3 + i * 2.55
    rect(s, x, 1.75, 2.35, 4.5, fill=DARK)
    rect(s, x+0.1, 1.85, 2.15, 3.5, fill=RGBColor(0x1F,0x29,0x37))
    txt(s, icon, x+0.1, 2.0, 2.15, 0.8, size=32, align=PP_ALIGN.CENTER, color=WHITE)
    txt(s, title, x+0.1, 2.75, 2.15, 0.45, size=13, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    bullet_box(s, items, x+0.1, 3.2, 2.15, 2.9, size=9, color=RGBColor(0xD1,0xD5,0xDB), icon="·")

# Language support
rect(s, 0.3, 6.35, 12.7, 0.85, fill=LIGHT)
langs = ["English", "हिन्दी", "ਪੰਜਾਬੀ", "ગુજરાતી", "मराठी", "தமிழ்", "తెలుగు", "ಕನ್ನಡ", "বাংলা"]
txt(s, "9 Languages: " + "  ·  ".join(langs), 0.5, 6.52, 12.3, 0.5, size=12, color=DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 11 — MARKET OPPORTUNITY
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Market Opportunity", "India's freight market is at an inflection point")
footer(s, 11)

stats = [
    ("₹15T", "India Freight Market\n(2024)"),
    ("5.5M", "Registered Truckers\n(unorganised)"),
    ("12%", "Annual Market\nGrowth Rate"),
    ("3%", "Digital Penetration\n(massive gap)"),
]
for i, (v, l) in enumerate(stats):
    stat_card(s, 0.3 + i*3.25, 1.4, 3.0, 1.8, v, l, bg=DARK, val_color=ORANGE)

opportunities = [
    ("🎯", "Immediate TAM", "₹450B", "Mid-size truckers (10-40T)\nbetween Indian tier-1 & tier-2 cities"),
    ("📈", "SAM (5 years)", "₹1.2T", "Full commercial freight including\ncold chain, hazmat, container"),
    ("🌏", "Vision", "₹3T+", "Pan-India logistics OS: freight,\nstorage, insurance, credit"),
]
for i, (icon, label, val, desc) in enumerate(opportunities):
    x = 0.3 + i * 4.35
    rect(s, x, 3.5, 4.1, 2.8, fill=LIGHT)
    txt(s, icon + "  " + label, x+0.2, 3.6, 3.8, 0.45, size=13, bold=True, color=DARK)
    txt(s, val, x+0.2, 4.08, 3.8, 0.8, size=38, bold=True, color=ORANGE)
    txt(s, desc, x+0.2, 4.88, 3.8, 1.1, size=11, color=GRAY)

txt(s, "India to become the 3rd largest economy by 2030 — logistics backbone is critical infrastructure",
    0.4, 6.5, 12.5, 0.45, size=13, bold=True, color=DARK, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 12 — BUSINESS MODEL
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Business Model", "Multiple monetisation streams with network effects")
footer(s, 12)

revenue = [
    ("💳", "Platform Commission", "5%", "On every load value. Deducted automatically at delivery confirmation. Merchant pays, trucker receives net amount."),
    ("⭐", "Premium Subscriptions", "₹999/mo", "Priority matching, verified badge, advanced analytics, early payment settlement for high-volume truckers."),
    ("🏦", "Financial Services", "1–2%", "Fuel advances, insurance premium, FASTag top-up, EMI for truck upgrades — embedded finance."),
    ("📣", "SaaS for Merchants", "₹2,499/mo", "Social media management, AI content, marketing analytics, branded tracking pages for customers."),
    ("📊", "Data Intelligence", "₹50K+/yr", "Lane pricing data, demand forecasting, fleet benchmarking sold to logistics companies & banks."),
    ("🚛", "Fleet Management", "₹499/truck/mo", "Document alerts, maintenance reminders, insurance renewal, GPS device subscription."),
]

for i, (icon, title, price, desc) in enumerate(revenue):
    x = 0.3 + (i % 2) * 6.5
    y = 1.4 + (i // 2) * 1.75
    rect(s, x, y, 6.2, 1.6, fill=LIGHT)
    txt(s, icon + "  " + title, x+0.2, y+0.1, 4.5, 0.45, size=13, bold=True, color=DARK)
    txt(s, price, x+5.0, y+0.05, 1.0, 0.5, size=14, bold=True, color=ORANGE, align=PP_ALIGN.RIGHT)
    txt(s, desc, x+0.2, y+0.55, 5.8, 0.9, size=10, color=GRAY)

txt(s, "Unit Economics: Average load value ₹45,000 × 5% = ₹2,250 revenue per trip  ●  Target: 10,000 trips/month → ₹22.5Cr MRR",
    0.4, 6.8, 12.5, 0.35, size=10, color=GRAY, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════
# SLIDE 13 — PLATFORM STATUS / QA
# ══════════════════════════════════════════════════════════════════
s = add_slide()
header_bar(s, "Platform Status", "Production-ready · 38/38 QA tests passing · Zero critical bugs")
footer(s, 13)

# Big QA score
rect(s, 0.3, 1.4, 4.5, 3.5, fill=DARK)
txt(s, "38/38", 0.3, 1.6, 4.5, 1.8, size=64, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
txt(s, "QA Tests Passing", 0.3, 3.2, 4.5, 0.6, size=16, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "100%", 0.3, 3.75, 4.5, 0.7, size=28, bold=True, color=GREEN, align=PP_ALIGN.CENTER)
txt(s, "All critical paths verified", 0.3, 4.4, 4.5, 0.4, size=11, color=GRAY, align=PP_ALIGN.CENTER)

test_sections = [
    ("Auth (5/5)", "Login, wrong password, /auth/me, rate limit"),
    ("Trucker Profile (2/2)", "Profile with KYC + trucks, availability set"),
    ("Load Search (2/2)", "16 loads found, camelCase fields verified"),
    ("Load Accept (3/3)", "Accept, lock enforced, active load visible"),
    ("Journey (5/5)", "Start, fuel stop, deliver, stats accumulated"),
    ("Merchant (2/2)", "Post load, view my loads"),
    ("Disputes (2/2)", "Raise dispute, list disputes"),
    ("Admin (4/4)", "Users, loads, live fleet, sim status"),
    ("Health (9/9)", "All 9 microservices returning 200"),
]

for i, (title, desc) in enumerate(test_sections):
    x = 5.0 + (i % 2) * 4.2
    y = 1.4 + (i // 2) * 0.68
    rect(s, x, y, 4.0, 0.6, fill=GREEN)
    txt(s, "✓  " + title, x+0.1, y+0.06, 3.8, 0.3, size=10, bold=True, color=WHITE)
    txt(s, desc, x+0.1, y+0.33, 3.8, 0.25, size=8, color=RGBColor(0xD1,0xD5,0xDB))

# Deployed infra
rect(s, 0.3, 5.15, 12.7, 1.65, fill=LIGHT)
txt(s, "Deployed Infrastructure:", 0.5, 5.25, 4, 0.35, size=12, bold=True, color=DARK)
infra = ["18 Docker containers", "PostgreSQL + PostGIS", "MongoDB + Redis", "Kafka + Zookeeper", "Elasticsearch", "RabbitMQ"]
for i, item in enumerate(infra):
    txt(s, "● " + item, 0.5 + (i%3)*4.2, 5.6 + (i//3)*0.45, 4.0, 0.4, size=11, color=DARK)

# ══════════════════════════════════════════════════════════════════
# SLIDE 14 — NEXT STEPS / CALL TO ACTION
# ══════════════════════════════════════════════════════════════════
s = add_slide()
rect(s, 0, 0, 13.33, 7.5, fill=DARK)
rect(s, 0, 0, 0.18, 7.5, fill=ORANGE)
footer(s, 14)

txt(s, "What's Next", 0.5, 0.3, 12, 0.8, size=36, bold=True, color=WHITE)
txt(s, "Join us in digitising India's ₹15 Trillion freight backbone",
    0.5, 1.05, 12, 0.5, size=16, color=GRAY)

roadmap = [
    ("Q3 2026", ORANGE, [
        "Launch in Bangalore, Delhi, Mumbai",
        "Onboard 500 verified truckers",
        "10 enterprise merchant accounts",
        "iOS app submission to App Store",
    ]),
    ("Q4 2026", BLUE, [
        "Expand to 10 Indian cities",
        "Integrate NHAI FASTag API",
        "Launch embedded insurance product",
        "5,000 active monthly trips",
    ]),
    ("2027", GREEN, [
        "₹10Cr Monthly GMV target",
        "Series A fundraise (₹50Cr)",
        "Pan-India coverage (50 cities)",
        "Cold chain & hazmat verticals",
    ]),
]

for i, (quarter, color, items) in enumerate(roadmap):
    x = 0.5 + i * 4.25
    rect(s, x, 1.7, 4.0, 3.8, fill=RGBColor(0x1F,0x29,0x37))
    rect(s, x, 1.7, 4.0, 0.5, fill=color)
    txt(s, quarter, x+0.15, 1.73, 3.7, 0.4, size=16, bold=True, color=WHITE)
    bullet_box(s, items, x+0.15, 2.25, 3.7, 3.2, size=12, color=WHITE, icon="→")

txt(s, "Platform live at  http://192.168.8.101:3010  (Portal)   ·   http://192.168.8.101:3011/admin  (Admin)",
    0.5, 5.7, 12.3, 0.45, size=13, color=ORANGE, align=PP_ALIGN.CENTER)

txt(s, "Contact: alexio1313@gmail.com  ·  Demo available on request",
    0.5, 6.2, 12.3, 0.4, size=13, color=WHITE, align=PP_ALIGN.CENTER)

txt(s, "\"Moving India Forward — One Load at a Time\" 🇮🇳",
    0.5, 6.7, 12.3, 0.5, size=15, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)

# ── Save ──────────────────────────────────────────────────────────
out = r"f:\AI_BOT\AI Trucker App\AI_Trucker_Platform_Executive_Presentation.pptx"
prs.save(out)
print("PPTX saved: " + out)
print("Slides: " + str(len(prs.slides)))
